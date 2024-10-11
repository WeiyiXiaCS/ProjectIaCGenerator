import boto3
import json
import base64
import time
import os
from pptx import Presentation
from docx import Document
import openpyxl
from io import BytesIO

def handler(event, context):
    bucket_name = "dupbucketforuploadfilesdemofrontendweiyi"
    key = event['Records'][0]['s3']['object']['key']
    file_url = f"s3://{bucket_name}/{key}"
    
    s3 = boto3.client('s3')
    textract = boto3.client('textract', region_name='eu-west-1')
    bedrock = boto3.client('bedrock-runtime', region_name='us-east-1')
    dynamodb = boto3.resource('dynamodb')
    table = dynamodb.Table('SyntheseDemoFrontend')
    
    try:
        time.sleep(5)  # Allow time for the uploading process to complete
        response = s3.get_object(Bucket=bucket_name, Key=key)
        file_content = response['Body'].read()
        file_extension = os.path.splitext(key)[1].lower()

        if file_extension in ['.jpeg', '.jpg', '.png']:
            content = get_content(file_content, file_extension, True)
        elif file_extension == '.pdf':
            extracted_text = extract_text_with_textract(textract, bucket_name, key)
            content = get_content(extracted_text, file_extension, False)
        elif file_extension == '.txt':
            content = get_content(file_content.decode('utf-8'), file_extension, False)
        elif file_extension in ['.ppt', '.pptx', '.doc', '.docx', '.xls', '.xlsx']:
            extracted_text = extract_text_from_office_file(file_content, file_extension)
            content = get_content(extracted_text, file_extension, False)
        else:
            raise Exception('Unsupported file type for text extraction.')
        
        prompt_message = (
            "The uploaded document is a resource for a project that needs to be preprocessed to retain only the core elements. "
            "Please provide: First, the essential contents of the project described in the document; Second, a thesis summarizing the document. "
            "Return the response in English in a strict JSON format without any additional words. "
            "The JSON format should be: {\"analysis\": \"essential contents\", \"thesis\": \"summary of the document\"}"
        )

        content.append({"type": "text", "text": prompt_message})

        # Invoke model with content
        result_json = invoke_model(bedrock, content)
        analysis, thesis = parse_bedrock_response(result_json)

        # Store the results to DynamoDB
        store_to_dynamodb(table, key, analysis, thesis, file_url)

        return {
            'statusCode': 200,
            'body': json.dumps({'Message': 'Content processed and stored successfully'})
        }
        
    except Exception as e:
        print(e)
        return {
            'statusCode': 500,
            'error': str(e)
        }

def get_content(data, file_extension, is_image):
    if is_image:
        encoded_content = base64.b64encode(data).decode()
        mime_type = 'image/jpeg' if file_extension in ['.jpeg', '.jpg'] else 'image/png'
        return [
            {"type": "image", "source": {"type": "base64", "media_type": mime_type, "data": encoded_content}}
        ]
    else:
        return [
            {"type": "text", "text": data}
        ]

def extract_text_with_textract(textract_client, bucket_name, key):
    job_id = start_pdf_job(textract_client, bucket_name, key)
    print("Started PDF job with id: {}".format(job_id))
    
    if is_job_complete(textract_client, job_id):
        response_pages = get_pdf_job_results(textract_client, job_id)
        extracted_text = ' '.join([item['Text'] for page in response_pages for item in page['Blocks'] if item['BlockType'] == 'LINE'])
        return extracted_text
    else:
        raise Exception("Textract processing did not complete successfully")

def start_pdf_job(client, s3_bucket_name, object_name):
    response = client.start_document_text_detection(
        DocumentLocation={'S3Object': {'Bucket': s3_bucket_name, 'Name': object_name}})
    return response["JobId"]

def is_job_complete(client, job_id):
    time.sleep(1)  # Initial sleep to allow job to start
    response = client.get_document_text_detection(JobId=job_id)
    status = response["JobStatus"]
    print("Job status: {}".format(status))

    while status == "IN_PROGRESS":
        time.sleep(1)
        response = client.get_document_text_detection(JobId=job_id)
        status = response["JobStatus"]
        print("Job status: {}".format(status))

    return status == "SUCCEEDED"

def get_pdf_job_results(client, job_id):
    pages = []
    next_token = None
    response = client.get_document_text_detection(JobId=job_id)
    if 'Blocks' in response:
        pages.append(response)

    next_token = response.get('NextToken', None)
    while next_token:
        time.sleep(1)  # Sleep between API calls to avoid rate limits
        response = client.get_document_text_detection(JobId=job_id, NextToken=next_token)
        pages.append(response)
        next_token = response.get('NextToken', None)

    return pages

def invoke_model(bedrock, content):
    body = json.dumps({
        "anthropic_version": "bedrock-2023-05-31",
        "max_tokens": 1000,
        "messages": [{"role": "user", "content": content}]
    })

    modelId = 'anthropic.claude-3-sonnet-20240229-v1:0'
    accept = 'application/json'
    contentType = 'application/json'

    try:
        response = bedrock.invoke_model(body=body, modelId=modelId, accept=accept, contentType=contentType)
        response_body = response['body'].read().decode('utf-8')
        result = json.loads(response_body)
        return result["content"][0]["text"]
    except Exception as e:
        print(e)
        return None

def parse_bedrock_response(response):
    try:
        result = json.loads(response)
        print(result)
        analysis = result.get("analysis","")
        thesis = result.get("thesis", "")
        return analysis, thesis
    except json.JSONDecodeError as e:
        print(f"Error parsing Bedrock response: {e}")
        return "", ""

def store_to_dynamodb(table, original_key, analysis, thesis, file_url):
    new_key = original_key
    try:
        table.put_item(Item={
            'Filename': new_key,
            'AdditionalInfo': "",
            'Analysis': analysis,
            'Thesis': thesis,
            'S3_URL': file_url
        })
        return True
    except Exception as e:
        print(e)
        return False

def extract_text_from_office_file(file_content, file_extension):
    file_stream = BytesIO(file_content)

    if file_extension in ['.ppt', '.pptx']:
        prs = Presentation(file_stream)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    text += shape.text + "\n"
        return text
    
    elif file_extension in ['.doc', '.docx']:
        doc = Document(file_stream)
        return "\n".join([paragraph.text for paragraph in doc.paragraphs])
    
    elif file_extension in ['.xls', '.xlsx']:
        wb = openpyxl.load_workbook(file_stream)
        text = ""
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            for row in ws.iter_rows(values_only=True):
                text += " ".join([str(cell) for cell in row if cell is not None]) + "\n"
        return text
    
    else:
        raise ValueError(f"Unsupported file type: {file_extension}")

